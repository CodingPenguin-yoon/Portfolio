from __future__ import annotations

from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
OUT = ROOT
SLIDE_OUT = ROOT / "pptx-slides"

W = 13.333
H = 7.5

INK = RGBColor(21, 21, 21)
PAPER = RGBColor(248, 248, 241)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(88, 96, 104)
GREEN = RGBColor(47, 157, 98)
BLUE = RGBColor(39, 86, 216)
RED = RGBColor(217, 74, 58)
YELLOW = RGBColor(231, 174, 53)
DEEP = RGBColor(16, 32, 25)
LINE = RGBColor(216, 222, 212)

FONT = "Apple SD Gothic Neo"
FONT_EN = "Inter"


def inch(value: float):
    return Inches(value)


def rgb(color: RGBColor):
    return color


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = inch(W)
    prs.slide_height = inch(H)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=PAPER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inch(W), inch(H))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def stripe(slide, horizontal: bool = True):
    colors = [GREEN, BLUE, RED, YELLOW]
    if horizontal:
        segment = W / 4
        for idx, color in enumerate(colors):
            rect(slide, segment * idx, H - 0.1, segment + 0.01, 0.1, color, line=False)
    else:
        segment = H / 4
        for idx, color in enumerate(colors):
            rect(slide, 0, segment * idx, 0.16, segment + 0.01, color, line=False)


def rect(slide, x, y, w, h, color, line=True, line_color=INK, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def text_box(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = FONT,
    line_spacing: float | None = None,
):
    shape = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def label(slide, text: str, x: float, y: float, w: float = 5.0, color: RGBColor = GREEN):
    return text_box(slide, text, x, y, w, 0.28, 12, color, True, font=FONT_EN)


def pill(slide, text: str, x: float, y: float, w: float):
    rect(slide, x, y, w, 0.33, WHITE, line=True)
    return text_box(slide, text, x + 0.07, y + 0.065, w - 0.14, 0.16, 9, INK, True, PP_ALIGN.CENTER, FONT_EN)


def shadow_card(slide, x, y, w, h, color=WHITE, shadow=True):
    if shadow:
        rect(slide, x + 0.07, y + 0.07, w, h, INK, line=False)
    return rect(slide, x, y, w, h, color, line=True)


def add_footer(slide, left: str, right: str):
    text_box(slide, left, 0.7, 7.03, 5.6, 0.2, 9, MUTED, True, font=FONT_EN)
    text_box(slide, right, 6.2, 7.03, 6.4, 0.2, 9, MUTED, True, PP_ALIGN.RIGHT, FONT_EN)


def add_code_bars(slide):
    # A quiet interface hint used behind title slides.
    for i, width in enumerate([1.2, 2.8, 2.0, 3.2, 1.7, 2.5]):
        rect(slide, 8.8, 1.0 + i * 0.48, width, 0.13, LINE, line=False, radius=True)


def slide_cover(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    image = ASSETS / "klepaas-dashboard.png"
    if image.exists():
        slide.shapes.add_picture(str(image), inch(3.8), 0, width=inch(9.7), height=inch(5.3))
        overlay = rect(slide, 0, 0, W, H, PAPER, line=False)
        overlay.fill.fore_color.transparency = 18
    stripe(slide, horizontal=False)
    label(slide, "PLATFORM OPERATIONS PORTFOLIO", 0.55, 0.45, 4.6)
    text_box(slide, "반복되는 운영을\n줄이고 싶었습니다.", 0.55, 1.6, 7.4, 1.45, 44, INK, True, line_spacing=0.9)
    text_box(
        slide,
        "처음엔 kubectl, 콘솔, 알림을 왔다 갔다 하는 게 번거로웠습니다. 그래서 자주 하는 일을 묶고, 상태가 남고, 되돌릴 수 있는 도구를 만들었습니다.",
        0.58,
        4.05,
        7.1,
        1.12,
        17,
        RGBColor(48, 54, 59),
        True,
    )
    text_box(slide, "Cho YunHo", 0.6, 6.35, 2.4, 0.28, 20, INK, True, font=FONT_EN)
    text_box(slide, "운영 자동화 · 플랫폼 엔지니어링", 0.6, 6.75, 3.6, 0.22, 12, MUTED, True)
    x = 8.25
    for name, width in [("FastAPI", 0.85), ("Kubernetes", 1.18), ("Docker", 0.82), ("Proxmox", 0.94)]:
        pill(slide, name, x, 6.55, width)
        x += width + 0.12
    return slide


def slide_story(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    stripe(slide)
    label(slide, "STORY LINE", 0.6, 0.55, 2.0)
    text_box(slide, "출발점은 단순했습니다.", 0.6, 1.05, 8.8, 0.55, 36, INK, True)
    text_box(
        slide,
        "같은 작업을 두 번 하고 싶지 않았고, 실행한 일은 나중에 확인할 수 있어야 한다고 생각했습니다.",
        0.62,
        1.82,
        8.9,
        0.42,
        18,
        RGBColor(48, 54, 59),
        True,
    )
    cards = [
        ("01", "K-Le-PaaS", "배포, 롤백, 알림, 헬스체크를 한 곳에서 처리하려고 만들었습니다.", RGBColor(231, 244, 223)),
        ("02", "Heimdall", "브랜치가 바뀌면 미리보기 환경이 뜨고, 로그와 롤백이 남게 했습니다.", RGBColor(232, 237, 255)),
        ("03", "Gjallar", "VM 운영은 실수 비용이 커서, 추천과 실행을 분리하고 승인 단계를 뒀습니다.", RGBColor(255, 233, 228)),
    ]
    for idx, (num, title, body, color) in enumerate(cards):
        x = 0.7 + idx * 4.1
        shadow_card(slide, x, 3.0, 3.35, 2.85, color)
        text_box(slide, num, x + 0.18, 3.22, 1.0, 0.42, 30, RGBColor(100, 105, 110), True, font=FONT_EN)
        text_box(slide, title, x + 0.18, 4.08, 2.7, 0.3, 20, INK, True, font=FONT_EN)
        text_box(slide, body, x + 0.18, 4.55, 2.9, 0.75, 12, RGBColor(57, 64, 71), True)
    return slide


def slide_klepaas(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    stripe(slide)
    label(slide, "PROJECT 01 · K-LE-PAAS", 0.6, 0.55, 3.2)
    text_box(slide, "kubectl과 콘솔을\n오가는 일을 줄이고\n싶었습니다.", 0.6, 1.02, 5.0, 1.6, 31, INK, True, line_spacing=0.9)
    text_box(
        slide,
        "K-Le-PaaS는 GitHub, Kubernetes, NCP, Slack, 모니터링을 FastAPI 백엔드로 묶은 프로젝트입니다. 배포하고, 롤백하고, 상태를 보고, 알림을 받는 일을 한 흐름으로 줄였습니다.",
        0.62,
        2.85,
        4.95,
        1.05,
        13,
        RGBColor(48, 54, 59),
        True,
    )
    points = [
        ("불편", "명령어, 콘솔, 알림이 따로 있어서 매번 확인할 곳이 많았습니다."),
        ("구현", "자연어 명령, Webhook, Kubernetes/NCP 작업, Slack 알림을 백엔드로 묶었습니다."),
        ("결과", "누가 무엇을 실행했는지 남고, 같은 운영 작업을 더 짧게 처리할 수 있게 했습니다."),
    ]
    for idx, (head, body) in enumerate(points):
        y = 3.95 + idx * 0.72
        shadow_card(slide, 0.65, y, 4.85, 0.54, WHITE, shadow=False)
        text_box(slide, head, 0.82, y + 0.13, 0.6, 0.16, 9, GREEN, True)
        text_box(slide, body, 1.45, y + 0.11, 3.75, 0.22, 9, RGBColor(57, 64, 71), True)
    image = ASSETS / "klepaas-dashboard.png"
    if image.exists():
        shadow_card(slide, 6.1, 1.25, 6.55, 4.45, WHITE)
        slide.shapes.add_picture(str(image), inch(6.1), inch(1.25), width=inch(6.55), height=inch(3.69))
        text_box(slide, "K-Le-PaaS Dashboard", 6.25, 5.12, 3.0, 0.18, 10, MUTED, True, font=FONT_EN)
    add_footer(slide, "github.com/K-Le-PaaS/backend-hybrid", "FastAPI · Kubernetes · NCP · Slack · Gemini NLP")
    return slide


def slide_heimdall(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    stripe(slide)
    label(slide, "PROJECT 02 · HEIMDALL", 0.6, 0.55, 3.0)
    text_box(slide, "푸시한 뒤\n“어디서 확인하지?”를\n없애고 싶었습니다.", 0.6, 1.02, 8.4, 1.25, 31, INK, True, line_spacing=0.9)
    text_box(
        slide,
        "Heimdall은 GitHub/GitLab 프로젝트를 등록해두면 브랜치 변경을 감지하고, Docker 기반 미리보기 환경을 띄워주는 배포 매니저입니다. 상태, 로그, 히스토리, 롤백까지 한 화면에서 확인하게 했습니다.",
        0.62,
        2.55,
        10.7,
        0.72,
        14,
        RGBColor(48, 54, 59),
        True,
    )
    steps = [
        ("1", "Register", "프로젝트와 브랜치를 한 번 등록"),
        ("2", "Build", "Webhook을 받으면 Docker 이미지 빌드"),
        ("3", "Preview", "URL, 상태, 로그를 바로 확인"),
        ("4", "Rollback", "문제 있으면 이전 이미지로 되돌림"),
    ]
    for idx, (num, title, body) in enumerate(steps):
        x = 0.72 + idx * 3.08
        shadow_card(slide, x, 3.35, 2.45, 1.72, WHITE)
        rect(slide, x + 0.18, 3.58, 0.36, 0.36, YELLOW)
        text_box(slide, num, x + 0.18, 3.63, 0.36, 0.1, 10, INK, True, PP_ALIGN.CENTER, FONT_EN)
        text_box(slide, title, x + 0.18, 4.18, 1.8, 0.22, 15, INK, True, font=FONT_EN)
        text_box(slide, body, x + 0.18, 4.52, 2.0, 0.35, 10, RGBColor(57, 64, 71), True)
    rect(slide, 0.75, 5.55, 11.75, 0.74, DEEP, line=False)
    text_box(slide, "제품 관점", 1.0, 5.78, 1.0, 0.15, 11, WHITE, True)
    text_box(slide, "코드를 올린 뒤 결과를 따로 찾아다니지 않고, 지금 어떤 버전이 떠 있는지와 되돌릴 수 있는지를 바로 보게 만들었습니다.", 2.18, 5.75, 9.7, 0.22, 11, RGBColor(225, 231, 226), True)
    add_footer(slide, "github.com/CodingPenguin-yoon/Heimdall", "FastAPI · React/Vite · Docker · PostgreSQL · Webhook")
    return slide


def slide_gjallar(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    stripe(slide)
    label(slide, "PROJECT 03 · GJALLAR", 0.6, 0.55, 2.8)
    text_box(slide, "위험한 작업은\n한 번 더 멈추게\n하고 싶었습니다.", 0.6, 1.02, 8.8, 1.25, 31, INK, True, line_spacing=0.9)
    text_box(
        slide,
        "Gjallar는 Proxmox 운영을 다루려고 만든 콘솔입니다. VM 이동 같은 작업은 자동으로 바로 실행하기보다, 현재 상태를 보고 추천을 만들고, 승인한 뒤에만 실행되게 나눴습니다.",
        0.62,
        2.55,
        10.6,
        0.72,
        14,
        RGBColor(48, 54, 59),
        True,
    )
    controls = [
        ("Check", "상태 확인", "VM, node, storage, HA, task 상태를 Proxmox에서 먼저 읽습니다.", GREEN),
        ("Suggest", "추천 만들기", "CPU/Memory, 정책, 리스크를 보고 이동해도 되는지 판단합니다.", YELLOW),
        ("Approve", "승인 후 실행", "실행은 승인과 live check를 통과한 job에서만 진행합니다.", RED),
    ]
    for idx, (small, title, body, color) in enumerate(controls):
        x = 0.75 + idx * 4.05
        shadow_card(slide, x, 3.35, 3.35, 1.95, WHITE)
        rect(slide, x + 0.18, 3.58, 0.9, 0.25, color, line=False)
        text_box(slide, small, x + 0.22, 3.635, 0.82, 0.08, 8, WHITE if color != YELLOW else INK, True, PP_ALIGN.CENTER, FONT_EN)
        text_box(slide, title, x + 0.18, 4.18, 2.4, 0.24, 16, INK, True, font=FONT_EN)
        text_box(slide, body, x + 0.18, 4.55, 2.85, 0.42, 10, RGBColor(57, 64, 71), True)
    rect(slide, 0.75, 5.78, 11.75, 0.6, DEEP, line=False)
    text_box(slide, "제품 관점", 1.0, 5.97, 1.0, 0.15, 11, WHITE, True)
    text_box(slide, "자동으로 다 해주는 것보다, 실수하면 큰 작업은 사람이 확인하고 승인할 수 있게 만드는 쪽에 집중했습니다.", 2.18, 5.94, 9.4, 0.2, 11, RGBColor(225, 231, 226), True)
    add_footer(slide, "github.com/CodingPenguin-yoon/Gjallar", "FastAPI · React · Proxmox API · Alembic · PostgreSQL · Audit")
    return slide


def slide_capability(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    stripe(slide)
    label(slide, "CAPABILITY SUMMARY", 0.6, 0.55, 3.2)
    text_box(slide, "제가 만든 것들의 공통점은\n반복 작업을 줄이는 것입니다.", 0.6, 1.05, 10.2, 0.95, 34, INK, True, line_spacing=0.9)
    items = [
        ("자주 하는 일 묶기", "배포, 롤백, 알림, 상태 확인처럼 반복되는 일을 한 화면과 API로 묶습니다.", GREEN),
        ("결과 남기기", "실행한 일의 로그, 히스토리, 현재 상태를 남겨 나중에 다시 볼 수 있게 합니다.", BLUE),
        ("위험한 일 멈추기", "VM 이동처럼 실수 비용이 큰 일은 바로 실행하지 않고 확인과 승인을 거치게 합니다.", RED),
        ("도구 사이 간격 줄이기", "kubectl, Proxmox, NCP, Slack처럼 떨어진 도구를 사용 흐름에 맞게 연결합니다.", YELLOW),
    ]
    for idx, (title, body, color) in enumerate(items):
        x = 0.85 + (idx % 2) * 5.85
        y = 2.75 + (idx // 2) * 1.72
        shadow_card(slide, x, y, 5.15, 1.18, WHITE)
        rect(slide, x, y, 5.15, 0.08, color, line=False)
        text_box(slide, title, x + 0.22, y + 0.28, 2.5, 0.2, 16, INK, True)
        text_box(slide, body, x + 0.22, y + 0.64, 4.55, 0.32, 10, RGBColor(57, 64, 71), True)
    return slide


def slide_closing(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, DEEP)
    stripe(slide)
    label(slide, "CLOSING", 0.6, 0.55, 1.8, RGBColor(142, 224, 173))
    text_box(slide, "자동화는 결국\n사람이 덜 헷갈리게\n만드는 일이라고\n생각합니다.", 0.6, 1.1, 8.5, 2.6, 37, WHITE, True, line_spacing=0.88)
    text_box(
        slide,
        "K-Le-PaaS, Heimdall, Gjallar는 모두 같은 불편함에서 시작했습니다. 매번 반복하는 일을 줄이고, 결과를 남기고, 위험한 일은 한 번 더 확인하게 만드는 것. 그런 도구를 계속 만들고 있습니다.",
        0.62,
        4.62,
        9.2,
        0.8,
        16,
        RGBColor(224, 232, 226),
        True,
    )
    for idx, (head, body) in enumerate([("GitHub", "github.com/CodingPenguin-yoon"), ("Projects", "K-Le-PaaS · Heimdall · Gjallar")]):
        x = 0.75 + idx * 6.0
        rect(slide, x, 6.25, 5.1, 0.62, DEEP, line=True, line_color=RGBColor(210, 225, 216))
        text_box(slide, head.upper(), x + 0.18, 6.39, 0.9, 0.1, 8, RGBColor(142, 224, 173), True, font=FONT_EN)
        text_box(slide, body, x + 1.05, 6.36, 3.7, 0.14, 13, WHITE, True, font=FONT_EN)
    return slide


SLIDES: list[tuple[str, Callable[[Presentation], object]]] = [
    ("01-cover", slide_cover),
    ("02-story-line", slide_story),
    ("03-klepaas", slide_klepaas),
    ("04-heimdall", slide_heimdall),
    ("05-gjallar", slide_gjallar),
    ("06-capability", slide_capability),
    ("07-closing", slide_closing),
]


def build_full() -> Path:
    prs = new_deck()
    for _, builder in SLIDES:
        builder(prs)
    out = OUT / "cho-yunho-portfolio.pptx"
    prs.save(out)
    return out


def build_split() -> list[Path]:
    SLIDE_OUT.mkdir(exist_ok=True)
    paths: list[Path] = []
    for name, builder in SLIDES:
        prs = new_deck()
        builder(prs)
        out = SLIDE_OUT / f"{name}.pptx"
        prs.save(out)
        paths.append(out)
    return paths


def main():
    full = build_full()
    split = build_split()
    print(f"created {full.relative_to(ROOT.parent)}")
    for path in split:
        print(f"created {path.relative_to(ROOT.parent)}")


if __name__ == "__main__":
    main()
